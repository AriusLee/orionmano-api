import json
import os
import re
import fitz  # pymupdf

from app.services.ai.client import generate_text


IMAGE_EXTS = (".jpg", ".jpeg", ".png", ".gif", ".webp", ".heic", ".heif", ".bmp", ".tiff")


# Filename keyword → category map. Order matters: more specific keywords first so
# generic tokens (e.g. "incorporation" → legal) don't swallow narrower ones
# (e.g. "tin certificate" → tax_return). Filenames are normalized to lowercase
# alphanumeric tokens before matching, so keywords here should be too.
FILENAME_KEYWORDS: list[tuple[str, str]] = [
    # tax
    ("cp204", "tax_return"),
    ("lhdn", "tax_return"),
    ("tin certificate", "tax_return"),
    ("tin", "tax_return"),
    ("ea form", "tax_return"),
    ("tax return", "tax_return"),
    ("tax computation", "tax_return"),
    ("tax filing", "tax_return"),
    # audit
    ("audited", "audit_report"),
    ("audit report", "audit_report"),
    ("auditor", "audit_report"),
    ("afs", "audit_report"),  # Audited Financial Statements
    ("statutory audit", "audit_report"),
    # management accounts
    ("management accounts", "management_accounts"),
    ("management account", "management_accounts"),
    ("mgmt account", "management_accounts"),
    # financial statements (extracts from audited FS, typically used as valuation/DD inputs)
    # — distinct from audit_report (no auditor opinion / notes) and from management_accounts
    # (those are interim unaudited; this is compiled FS line-items)
    ("financial statements", "financial_statements"),
    ("fs for valuation", "financial_statements"),
    ("fs", "financial_statements"),
    ("financials", "financial_statements"),
    # prospectus
    ("prospectus", "prospectus"),
    ("offering memorandum", "prospectus"),
    # org chart
    ("org chart", "org_chart"),
    ("organization chart", "org_chart"),
    ("organisation chart", "org_chart"),
    ("corporate structure", "org_chart"),
    ("group structure", "org_chart"),
    ("holding structure", "org_chart"),
    # cap table
    ("cap table", "cap_table"),
    ("shareholder register", "cap_table"),
    ("share register", "cap_table"),
    # shareholder agreement
    ("shareholders agreement", "shareholder_agreement"),
    ("shareholder agreement", "shareholder_agreement"),
    ("investment agreement", "shareholder_agreement"),
    ("subscription agreement", "shareholder_agreement"),
    ("term sheet", "shareholder_agreement"),
    ("sha", "shareholder_agreement"),
    # board minutes
    ("board minutes", "board_minutes"),
    ("board resolution", "board_minutes"),
    ("directors resolution", "board_minutes"),
    ("minutes of meeting", "board_minutes"),
    ("written resolution", "board_minutes"),
    # material contracts
    ("licensing agreement", "material_contract"),
    ("franchise agreement", "material_contract"),
    ("distribution agreement", "material_contract"),
    ("supply agreement", "material_contract"),
    ("mou", "material_contract"),
    # projections
    ("projection", "projections"),
    ("forecast", "projections"),
    ("budget", "projections"),
    ("business plan", "projections"),
    ("financial model", "projections"),
    # company profile
    ("company profile", "company_profile"),
    ("pitch deck", "company_profile"),
    ("company overview", "company_profile"),
    # interviews
    ("interview", "interview"),
    # legal (generic keywords last — "certificate" and "incorporation" match many filenames)
    ("certificate of incorporation", "legal"),
    ("memorandum of association", "legal"),
    ("articles of association", "legal"),
    ("legal opinion", "legal"),
    ("litigation", "legal"),
    ("incorporation", "legal"),
    ("ssm", "legal"),
    ("companies act", "legal"),
    ("constitution", "legal"),
]


def classify_by_filename(filename: str | None) -> str | None:
    """Best-effort classifier for when LLM extraction yields no document_type
    (scanned PDFs, image uploads, parse errors). Returns a taxonomy slug or None."""
    if not filename:
        return None
    norm = " " + re.sub(r"[^a-z0-9]+", " ", filename.lower()).strip() + " "
    for keyword, category in FILENAME_KEYWORDS:
        if f" {keyword} " in norm:
            return category
    return None


def _reconcile_categories(
    llm_cats: list[str] | None,
    filename: str | None,
) -> tuple[list[str], str]:
    """Combine LLM-derived categories with the filename keyword classifier.

    Filename keywords are a curated, high-precision allowlist — when one fires,
    it beats a conflicting LLM choice. (LLM has known failure modes on file
    content alone, e.g. mistaking a DCF projection workpaper for management
    accounts because it sees Profit-before-tax line items.) The LLM's pick is
    preserved as a secondary entry so multi-slot rendering still works.

    Returns (categories, classification_source)."""
    llm_clean = [c.strip().lower() for c in (llm_cats or []) if isinstance(c, str) and c.strip()]
    # Dedup while preserving order
    seen: set[str] = set()
    llm_clean = [c for c in llm_clean if not (c in seen or seen.add(c))]

    guessed = classify_by_filename(filename)

    if guessed and guessed != "other":
        if not llm_clean or llm_clean[0] != guessed:
            secondaries = [c for c in llm_clean if c != guessed and c != "other"]
            return [guessed] + secondaries, "filename_priority"
        return llm_clean, "llm"

    if llm_clean and llm_clean != ["other"]:
        return llm_clean, "llm"

    return ["other"], "default"


def extract_text_from_pdf(file_path: str, max_pages: int = 400) -> str:
    # Cap is high (not 50) because prospectus / DRS F-pages — the audited
    # financial statements — sit near the BACK of the document. A low cap meant
    # the financials were never read, so financial_data came back empty. fitz
    # text extraction is cheap even for several-hundred-page filings; the real
    # bound on what reaches the LLM is _focused_excerpt() below.
    doc = fitz.open(file_path)
    text_parts = []
    for i, page in enumerate(doc):
        if i >= max_pages:
            break
        text_parts.append(page.get_text())
    doc.close()
    return "\n\n".join(text_parts)


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a .docx, including TABLE content in document order.

    A .docx is a ZIP archive, not text — the previous code path read it in text
    mode and got raw ZIP bytes ("PK\\x03\\x04…"), so Word prospectuses were never
    parsed (not classified, no financials). Financial statements in a prospectus
    live almost entirely in tables, so we walk the body and render both
    paragraphs and table rows (cells pipe-joined) in order."""
    from docx import Document as _Docx
    from docx.oxml.ns import qn

    document = _Docx(file_path)
    parts: list[str] = []
    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            line = "".join(t.text or "" for t in child.findall(".//" + qn("w:t")))
            if line.strip():
                parts.append(line)
        elif child.tag == qn("w:tbl"):
            for row in child.findall(qn("w:tr")):
                cells = [
                    "".join(t.text or "" for t in tc.findall(".//" + qn("w:t"))).strip()
                    for tc in row.findall(qn("w:tc"))
                ]
                if any(cells):
                    parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_text_from_xlsx(file_path: str, max_chars: int = 300_000) -> str:
    """Extract text from a .xlsx/.xlsm, one block per sheet, rows as pipe-joined
    cells. Same class of bug as .docx — a spreadsheet is a ZIP archive, so the
    old text-mode read produced garbage and financial models / management
    accounts / cap tables never got parsed.

    `data_only=True` returns the last cached cell values rather than formula
    strings, so a financial model yields its numbers (not '=SUM(...)').
    `read_only=True` streams large books cheaply; we stop at max_chars and let
    _focused_excerpt() do the final bound for the LLM."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    parts: list[str] = []
    size = 0
    try:
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if not any(c.strip() for c in cells):
                    continue
                # trim trailing empties so wide-but-sparse rows stay compact
                while cells and not cells[-1].strip():
                    cells.pop()
                line = " | ".join(cells)
                parts.append(line)
                size += len(line) + 1
                if size >= max_chars:
                    parts.append("[... spreadsheet truncated ...]")
                    return "\n".join(parts)
    finally:
        wb.close()
    return "\n".join(parts)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a .pptx (pitch deck / company profile), one block per
    slide, including shape text and TABLE rows. Same ZIP-archive bug class as
    .docx/.xlsx — must be parsed, not read as text."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {idx}")
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
            elif shape.has_text_frame:
                txt = shape.text_frame.text
                if txt and txt.strip():
                    parts.append(txt.strip())
    return "\n".join(parts)


# Statement-title / line-item anchors that mark the actual audited financial
# statements (the "F-pages"). IFRS-friendly (HK/SG issuers say "profit/loss for
# the year", "statements of financial position") plus US-GAAP wording. We use
# these to locate the financials inside a long filing so the extractor sees them
# even when they sit 70k+ chars deep. Deliberately excludes the narrative
# "report of independent registered…" mention, which appears early (Experts/
# Summary) and is a false anchor for the statements themselves.
_FIN_STATEMENT_ANCHORS = (
    "consolidated statements of operations",
    "consolidated statement of operations",
    "consolidated statements of profit or loss",
    "consolidated statements of comprehensive income",
    "consolidated balance sheet",
    "consolidated balance sheets",
    "consolidated statements of financial position",
    "consolidated statement of financial position",
    "consolidated statements of cash flows",
    "total current assets",
    "total assets",
)


def _focused_excerpt(
    text: str,
    head_chars: int = 10_000,
    window_chars: int = 80_000,
    total_cap: int = 95_000,
) -> str:
    """Bound what gets sent to the extraction LLM while guaranteeing the audited
    financial statements are included for long filings (prospectus / DRS).

    Short docs (<= total_cap) are returned whole. For long docs we keep the head
    (cover page / company info → classification + company_info) and splice in the
    region around the first real financial-statement anchor (→ income statement,
    balance sheet, cash flows). This is the extraction-time companion to the
    report-side F-page source hierarchy."""
    if len(text) <= total_cap:
        return text

    low = text.lower()
    fin_pos = None
    for anchor in _FIN_STATEMENT_ANCHORS:
        i = low.find(anchor)
        if i != -1:
            fin_pos = i if fin_pos is None else min(fin_pos, i)

    # No financials located, or they already fall inside the head window — just
    # send the head-capped slice.
    if fin_pos is None or fin_pos < head_chars:
        return text[:total_cap]

    start = max(fin_pos - 3_000, head_chars)
    combined = (
        text[:head_chars]
        + "\n\n[... intervening narrative sections omitted for extraction ...]\n\n"
        + text[start:start + window_chars]
    )
    return combined[:total_cap]


async def extract_document(file_path: str, filename: str | None = None) -> dict:
    lower = file_path.lower()
    fname = filename or os.path.basename(file_path)

    # Images: hand to the vision model. A single image can legitimately satisfy
    # multiple slots (e.g. an org chart that also shows the cap table).
    if lower.endswith(IMAGE_EXTS):
        from app.services.ai.vision import classify_image_file

        try:
            v = await classify_image_file(file_path)
            cats, source = _reconcile_categories(v.get("categories"), fname)
            return {
                "document_type": cats[0],
                "categories": cats,
                "classification_source": source if source != "llm" else "vision",
                "summary": v.get("summary", ""),
            }
        except Exception as e:
            guessed = classify_by_filename(fname)
            return {
                "document_type": guessed or "other",
                "categories": [guessed] if guessed else ["other"],
                "classification_source": "filename" if guessed else "default",
                "error": f"Vision failed: {e}",
            }

    if lower.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
    elif lower.endswith(".docx"):
        # .docx is a ZIP archive — must be parsed, not read as text. If parsing
        # fails for any reason, fall back to the raw read so the doc still gets
        # a filename-based classification instead of erroring the whole upload.
        try:
            text = extract_text_from_docx(file_path)
        except Exception:
            text = ""
    elif lower.endswith((".xlsx", ".xlsm")):
        # .xlsx/.xlsm are ZIP archives too — same fix as .docx.
        try:
            text = extract_text_from_xlsx(file_path)
        except Exception:
            text = ""
    elif lower.endswith(".pptx"):
        # .pptx pitch decks / company profiles — same ZIP-archive fix.
        try:
            text = extract_text_from_pptx(file_path)
        except Exception:
            text = ""
    else:
        with open(file_path, "r", errors="ignore") as f:
            text = f.read()[:50000]

    # Scanned PDFs / empty docs — no text layer. Try vision on the first pages
    # before falling back to filename keywords.
    if not text.strip():
        if lower.endswith(".pdf"):
            from app.services.ai.vision import classify_pdf_via_vision

            try:
                v = await classify_pdf_via_vision(file_path, max_pages=2)
                cats, source = _reconcile_categories(v.get("categories"), fname)
                return {
                    "document_type": cats[0],
                    "categories": cats,
                    "classification_source": source if source != "llm" else "vision_pdf",
                    "summary": v.get("summary", ""),
                    "raw_text": "",
                }
            except Exception:
                pass

        guessed = classify_by_filename(fname)
        return {
            "document_type": guessed or "other",
            "categories": [guessed] if guessed else ["other"],
            "classification_source": "filename" if guessed else "scan_needed",
            "error": "No text content extracted",
            "raw_text": "",
        }

    system_prompt = """You are a financial document analyst for Orionmano Assurance Services.
Extract structured data from the provided document. Return valid JSON with the following structure
(include only fields that are present in the document).

For `categories`, return an ARRAY of every slug that applies — a single document can satisfy
multiple slots. E.g. an annex inside an audit report that includes the shareholder register
matches BOTH "audit_report" AND "cap_table". Be generous but accurate: include a slug only
when the document's content clearly supports it. `document_type` should be the single primary
category (the first / most representative of the list).

Valid slugs for `document_type` and `categories`:
- audit_report — audited financial statements, auditor's opinion, PCAOB/MIA statutory audit reports
- management_accounts — interim/management P&L, balance sheet, unaudited financials
- financial_statements — historical P&L / balance sheet / cash flow extracts compiled from audited statements (no auditor opinion/notes), typically as inputs to valuation or DD modeling. Use this for files like "FS for valuation" or "Historical financials" that show line-item financials across multiple periods without being a full audit report or interim management accounts.
- tax_return — tax returns, tax filings, tax computations, CP204/LHDN filings
- org_chart — organization chart, corporate structure chart, group holding diagram
- cap_table — cap table, shareholder register, shareholding structure, share ledger
- board_minutes — board minutes, board resolutions, committee minutes, written resolutions
- shareholder_agreement — shareholders agreement, investment agreement, SHA, term sheet
- material_contract — customer / supplier / distribution / licensing / franchise contracts
- company_profile — company profile, pitch deck, corporate proposal, introduction slides
- projections — financial projections, budgets, forecasts, business plans with forward numbers
- legal — legal opinion, litigation report, regulatory correspondence, compliance letter
- prospectus — prospectus, offering memorandum, registration statement (S-1/F-1/20-F drafts)
- interview — management interview transcript, Q&A notes
- other — anything that does not clearly match the above

{
  "document_type": "audit_report|management_accounts|financial_statements|tax_return|org_chart|cap_table|board_minutes|shareholder_agreement|material_contract|company_profile|projections|legal|prospectus|interview|other",
  "categories": ["audit_report"],
  "company_info": {
    "name": "",
    "legal_name": "",
    "registration_number": "",
    "incorporation_date": "",
    "jurisdiction": "",
    "industry": "",
    "description": "",
    "website": ""
  },
  "financial_data": {
    "currency": "",
    "periods": ["FY2023", "FY2024"],
    "income_statement": {
      "revenue": {},
      "cost_of_revenue": {},
      "gross_profit": {},
      "operating_expenses": {},
      "finance_costs": {},
      "profit_before_tax": {},
      "taxation": {},
      "net_income": {}
    },
    "balance_sheet": {
      "total_assets": {},
      "total_liabilities": {},
      "total_equity": {},
      "current_assets": {},
      "current_liabilities": {},
      "cash": {}
    },
    "cash_flow": {
      "operating": {},
      "investing": {},
      "financing": {},
      "net_change": {}
    }
  },
  "shareholders": [
    {"name": "", "shares": 0, "percentage": 0}
  ],
  "key_personnel": [
    {"name": "", "title": "", "background": ""}
  ],
  "key_findings": [""],
  "summary": ""
}

Only include sections where you found relevant data. Keep values as numbers where possible.
For financial data, use the period as key (e.g. {"FY2023": 7522, "FY2024": 15291}).
"""

    # Focused excerpt instead of a blind first-30k slice: for long filings this
    # guarantees the audited financial statements (F-pages) reach the extractor.
    # max_tokens lifted to 8192 (DeepSeek ceiling) so multi-period financial
    # tables don't truncate the JSON and trip the parse-error fallback (which
    # silently drops financial_data).
    result = await generate_text(
        system_prompt=system_prompt,
        user_prompt=f"Extract structured data from this document:\n\n{_focused_excerpt(text)}",
        max_tokens=8192,
    )

    try:
        # Try to parse JSON from the response
        # Handle case where response has markdown code blocks
        cleaned = result.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1]
            cleaned = cleaned.rsplit("```", 1)[0]
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        guessed = classify_by_filename(fname)
        return {
            "raw_extraction": result,
            "parse_error": True,
            "document_type": guessed or "other",
            "categories": [guessed] if guessed else ["other"],
            "classification_source": "filename" if guessed else "default",
        }

    # Normalize categories: prefer the array the LLM returned, fall back to the
    # singular document_type, then filename heuristic. Keep document_type synced
    # to the first entry for any legacy callers.
    if isinstance(parsed, dict):
        raw_cats = parsed.get("categories")
        cats: list[str] = []
        if isinstance(raw_cats, list):
            for c in raw_cats:
                if isinstance(c, str) and c.strip():
                    c = c.strip().lower()
                    if c not in cats:
                        cats.append(c)

        doc_type = str(parsed.get("document_type") or "").strip().lower()

        if not cats and doc_type and doc_type != "other":
            cats = [doc_type]

        cats, source = _reconcile_categories(cats, fname)
        parsed["categories"] = cats
        parsed["document_type"] = cats[0]
        parsed["classification_source"] = source
    return parsed
